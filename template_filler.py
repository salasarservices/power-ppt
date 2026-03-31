"""
template_filler.py

For each page extracted from the input presentation, copies the brand template
slide (preserving logo, footer, geometric shapes, and all master styling), then:
  - replaces "TITLE GOES HERE" with the slide's title (XML-level, preserves formatting)
  - adds body text or tables in the white content area
"""

import copy
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


# ── Content area in the Salasar brand template (13.33" × 7.5" widescreen) ────
# The large white/blank zone sits below the title bar and above the footer.
CONTENT_LEFT   = Inches(0.75)
CONTENT_TOP    = Inches(1.3)
CONTENT_WIDTH  = Inches(11.8)
CONTENT_HEIGHT = Inches(5.1)

# Text that marks the title placeholder in the brand template slide
_TITLE_MARKER = "TITLE GOES HERE"

# DrawingML namespace
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
# Relationship namespace
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Relationship attributes that carry rIds pointing to embedded parts
_RID_ATTRS = [
    f"{{{_R_NS}}}embed",
    f"{{{_R_NS}}}id",
    f"{{{_R_NS}}}link",
    f"{{{_R_NS}}}href",
]

# Tags that form the structural header of a shape-tree (not copied content)
_STRUCT_TAGS = {qn("p:cNvGrpSpPr"), qn("p:grpSpPr")}


# ── Slide copy helpers ────────────────────────────────────────────────────────

def _copy_rels(src_slide, dst_slide):
    """
    Copy every relationship from src_slide to dst_slide.
    Returns {old_rId: new_rId} for any IDs that changed.
    Handles both embedded (internal) and linked (external) relationships
    so that logos stored as external URLs are preserved correctly.
    """
    rId_map = {}
    src_part = src_slide.part
    dst_part = dst_slide.part

    for rel in list(src_part.rels.values()):
        try:
            new_rId = dst_part.relate_to(
                rel._target, rel.reltype, is_external=rel.is_external
            )
            if new_rId != rel.rId:
                rId_map[rel.rId] = new_rId
        except Exception:
            pass

    return rId_map


def _remap_rids(element, rId_map):
    """Walk an XML element tree and rewrite any rId attributes using rId_map."""
    if not rId_map:
        return
    for el in element.iter():
        for attr in _RID_ATTRS:
            if attr in el.attrib and el.attrib[attr] in rId_map:
                el.attrib[attr] = rId_map[el.attrib[attr]]


def _copy_slide(prs, source_idx=0):
    """
    Append a full copy of the slide at source_idx to `prs`.
    Copies shapes, images, and all slide-level relationships.
    Returns the new slide.
    """
    source = prs.slides[source_idx]
    layout = source.slide_layout

    # Add a blank slide using the same layout (inherits slide master)
    new_slide = prs.slides.add_slide(layout)

    # ── Replace the new slide's shape tree with a copy of the source's ───────
    src_tree = source.shapes._spTree
    dst_tree = new_slide.shapes._spTree

    # Remove all non-structural children that add_slide inserted
    for el in list(dst_tree):
        if el.tag not in _STRUCT_TAGS:
            dst_tree.remove(el)

    # Deep-copy every non-structural element from source
    copied = [copy.deepcopy(el) for el in src_tree if el.tag not in _STRUCT_TAGS]

    # Copy relationships first so we can remap rIds in the copied XML
    rId_map = _copy_rels(source, new_slide)
    for el in copied:
        _remap_rids(el, rId_map)
        dst_tree.append(el)

    return new_slide


def _fix_shape_ids(slide):
    """
    After copying template shapes onto a new slide, every copied shape keeps
    its original id attribute (e.g. id="1", id="2" …).  When python-pptx then
    calls add_textbox / add_table it starts its own counter from 1 — producing
    duplicate <p:cNvPr id="…"> values.  PowerPoint treats duplicate IDs as
    corruption and prompts "file needs repair".

    This function walks every <p:cNvPr> element in the slide's shape tree and
    assigns a unique sequential id so the invariant is satisfied.
    """
    sp_tree = slide.shapes._spTree
    counter = 1
    for elem in sp_tree.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("id", str(counter))
            counter += 1


def _remove_slide_number_fields(slide):
    """
    Remove any shapes that contain an <a:fld type="slidenum"> field from the
    copied slide.  In the brand template these show as "1", "2", … on every
    output slide, which is unwanted noise.
    """
    sp_tree = slide.shapes._spTree
    to_remove = []
    for sp in list(sp_tree.findall(qn("p:sp"))):
        for fld in sp.findall(f".//{{{_A_NS}}}fld"):
            if "slidenum" in fld.get("type", "").lower():
                to_remove.append(sp)
                break
    for sp in to_remove:
        sp_tree.remove(sp)


def _remove_slide(prs, index):
    """
    Remove the slide at `index` from the presentation.

    Reads the relationship ID directly from the <p:sldId r:id="…"> XML
    attribute — more reliable than the previous object-identity search which
    could silently fail and leave a dangling relationship that makes PowerPoint
    report the file as needing repair.
    """
    xml_slides = prs.slides._sldIdLst
    sld_id_elem = xml_slides[index]

    # rId lives in the r: (relationship) namespace on the element itself
    r_id = sld_id_elem.get(qn("r:id"))

    # Remove the <p:sldId> element
    xml_slides.remove(sld_id_elem)

    # Drop the corresponding relationship from the presentation part
    if r_id:
        try:
            prs.part.drop_rel(r_id)
        except Exception:
            pass


# ── Per-slide content helpers ─────────────────────────────────────────────────

def _set_title(slide, title_text):
    """
    Replace the "TITLE GOES HERE" text in the brand template slide with
    `title_text`. Operates at the XML level so all surrounding formatting
    (font, colour, size, bold) is preserved.

    Falls back to setting the first title placeholder if the marker is not found.
    """
    marker_lower = _TITLE_MARKER.lower()

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if marker_lower not in shape.text_frame.text.lower():
            continue

        # Locate all <a:t> text nodes inside this shape
        sp_elem = shape._element
        t_nodes = sp_elem.findall(f".//{{{_A_NS}}}t")

        # Clear every run's text, then put the new title into the first node
        for t in t_nodes:
            t.text = ""
        if t_nodes:
            t_nodes[0].text = title_text
        return

    # Fallback — standard title placeholder (type 1)
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == 1:
                ph.text = title_text
                return
        except Exception:
            pass


def _add_content(slide, body_text, tables_for_slide, body_font="Poppins"):
    """
    Add tables (preferred) or a body text box to the slide's white content area.
    """
    if tables_for_slide:
        current_top = CONTENT_TOP

        for table_data in tables_for_slide:
            rows = table_data.get("rows", [])
            if not rows:
                continue

            n_rows = len(rows)
            n_cols = len(rows[0])
            if n_rows == 0 or n_cols == 0:
                continue

            row_h = Inches(0.34)
            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                CONTENT_LEFT, current_top,
                CONTENT_WIDTH, row_h * n_rows,
            )
            tbl = tbl_shape.table

            for r_idx, row_data in enumerate(rows):
                for c_idx in range(min(len(row_data), n_cols)):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.text = str(row_data[c_idx])
                    for para in cell.text_frame.paragraphs:
                        para.font.name = body_font
                        if r_idx == 0:
                            para.font.bold = True
                            para.font.size = Pt(11)
                            para.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            para.font.size = Pt(10)

                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0, 82, 156)

            current_top += row_h * n_rows + Inches(0.2)

    elif body_text and body_text.strip():
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = body_text.strip()

        for para in tf.paragraphs:
            para.font.name = body_font
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(30, 30, 30)


# ── Public API ────────────────────────────────────────────────────────────────

def fill_template_with_pages(template_bytes, pages, tables,
                             title_font="Poppins", body_font="Poppins"):
    """
    For every page in `pages`, copy the brand template slide (slide 0),
    update its title, add its content, then remove the original template slide.

    Parameters
    ----------
    template_bytes : bytes   — the brand template PPTX
    pages          : list[dict] — each has keys: title, body, slide_index
    tables         : list[dict] — each has keys: rows, slide_index
    title_font     : str
    body_font      : str

    Returns
    -------
    bytes — the finished PPTX
    """
    if not template_bytes:
        raise ValueError("Brand template is required.")

    prs = Presentation(io.BytesIO(template_bytes))

    if not prs.slides:
        raise ValueError("Brand template has no slides.")

    n_template_slides = len(prs.slides)   # remember how many to remove later

    for page in pages:
        slide_idx  = page.get("slide_index", 0)
        title_text = page.get("title") or f"Slide {slide_idx + 1}"
        body_text  = page.get("body", "")

        # Copy the brand template slide (always source_idx=0 because originals
        # are at the front and new slides are appended at the end)
        new_slide = _copy_slide(prs, source_idx=0)
        _remove_slide_number_fields(new_slide)

        _set_title(new_slide, title_text)

        tables_for_slide = [t for t in tables if t.get("slide_index") == slide_idx]
        _add_content(new_slide, body_text, tables_for_slide, body_font)

        # Re-number all shape IDs so there are no duplicates after the copy
        _fix_shape_ids(new_slide)

    # Remove the original template slides (they are still at index 0…n-1)
    for _ in range(n_template_slides):
        _remove_slide(prs, 0)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
