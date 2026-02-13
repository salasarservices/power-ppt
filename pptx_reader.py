import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_shapes(pptx_bytes):
    """
    Extract text shapes, tables, and images from a PowerPoint presentation.
    
    Args:
        pptx_bytes: Binary content of the PPTX file.
    
    Returns:
        List of dictionaries containing slide metadata including:
        - slide_index
        - title_text
        - body_text
        - shapes (including tables and images)
        - text_shapes
        - image_shapes
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides_meta = []

    for slide_idx, slide in enumerate(prs.slides):
        meta = {
            "slide_index": slide_idx,
            "title_text": "",
            "body_text": "",
            "text_shapes": [],
            "image_shapes": [],
            "shapes": []
        }

        title_text = ""
        body_text_parts = []

        for shape in slide.shapes:
            # Store the shape for table extraction
            meta["shapes"].append(shape)

            # Extract title from title placeholder
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:  # Title placeholder (PP_PLACEHOLDER.TITLE = 1)
                    if shape.has_text_frame:
                        title_text = shape.text_frame.text.strip()

            # Extract text from text boxes and content placeholders
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Skip if this is the title (already captured)
                    if text != title_text:
                        body_text_parts.append(text)
                    
                    meta["text_shapes"].append({
                        "text": text,
                        "shape_type": shape.shape_type
                    })

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    meta["image_shapes"].append({
                        "image_bytes": image_bytes,
                        "content_type": image.content_type
                    })
                except Exception as e:
                    print(f"Failed to extract image from slide {slide_idx}: {e}")

        # Set extracted title and body
        meta["title_text"] = title_text if title_text else f"Slide {slide_idx + 1}"
        meta["body_text"] = "\n\n".join(body_text_parts) if body_text_parts else ""

        slides_meta.append(meta)

    return slides_meta
